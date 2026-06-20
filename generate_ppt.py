"""
Generate academic PPT for IRS Anti-Jamming THz paper.
Equations are rendered as matplotlib mathtext images and embedded as pictures.
Run: python generate_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import math_helpers as mh

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY  = RGBColor(0x1A, 0x3A, 0x6C)
BLUE  = RGBColor(0x1F, 0x77, 0xB4)
LGRAY = RGBColor(0xF2, 0xF4, 0xF7)
DGRAY = RGBColor(0x44, 0x44, 0x44)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
RED   = RGBColor(0xD6, 0x27, 0x28)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(sl, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.width = lw
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line
    else: s.line.fill.background()
    return s

def tb(sl, text, l, t, w, h, fs=15, bold=False, color=DGRAY,
       align=PP_ALIGN.LEFT, italic=False):
    b = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    b.word_wrap = True
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return b

def add_eq(sl, latex, l, t, w, h, fs=20, bg='#F5F8FF', tc='#1A3A6C', fig_h=None):
    """Render a LaTeX string as a matplotlib image and place it on the slide."""
    fh = fig_h if fig_h else max(0.6, h * 1.1)
    buf = mh.render_eq(latex, fig_w=max(6.0, w * 1.2), fig_h=fh,
                       fontsize=fs, bg_color=bg, text_color=tc)
    sl.shapes.add_picture(buf, Inches(l), Inches(t), Inches(w), Inches(h))

def hdr(sl, title, sub=None):
    add_rect(sl, 0, 0, 13.33, 1.05, fill=NAVY)
    tb(sl, title, 0.3, 0.08, 12.5, 0.75, fs=26, bold=True, color=WHITE)
    if sub:
        tb(sl, sub, 0.3, 0.75, 12.5, 0.32, fs=13,
           color=RGBColor(0xBB, 0xCC, 0xEE))
    add_rect(sl, 0, 1.05, 13.33, 0.04, fill=BLUE)

def buls(sl, items, l, t, w, h, fs=14, color=DGRAY, title=None, tc=NAVY):
    b = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    b.word_wrap = True; tf = b.text_frame; tf.word_wrap = True; first = True
    if title:
        p = tf.paragraphs[0]; first = False; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = title
        r.font.size = Pt(fs+1); r.font.bold = True; r.font.color.rgb = tc
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = PP_ALIGN.LEFT; p.space_before = Pt(3)
        r = p.add_run()
        r.text = ("▸  " if not item.startswith("  ") else "    –  ") + item.lstrip()
        r.font.size = Pt(fs); r.font.color.rgb = color

def ch(sl, text, l, t, w, h=0.38, bg=NAVY, fg=WHITE, fs=13):
    add_rect(sl, l, t, w, h, fill=bg)
    tb(sl, text, l+0.05, t+0.03, w-0.1, h-0.06, fs=fs, bold=True,
       color=fg, align=PP_ALIGN.CENTER)

def tr(sl, cells, l, t, widths, h=0.32, bg=WHITE, fg=DGRAY, fs=12, bold=False):
    x = l
    for cell, w in zip(cells, widths):
        add_rect(sl, x, t, w, h, fill=bg,
                 line=RGBColor(0xCC,0xCC,0xCC), lw=Pt(0.5))
        tb(sl, str(cell), x+0.05, t+0.04, w-0.1, h-0.06,
           fs=fs, bold=bold, color=fg, align=PP_ALIGN.CENTER)
        x += w

# ── Equation strings (matplotlib mathtext) ────────────────────────────────────
EQ = {
  'received_signal':
    r'$y_k[m] = \mathbf{h}_{eff,k}^H[m]\,\mathbf{w}_k[m]\sqrt{P_k}\,s_k'
    r' + \!\sum_{i \neq k}\!\mathbf{h}_{eff,i}^H\mathbf{w}_i\sqrt{P_i}s_i'
    r' + \mathbf{h}_{J,k}^H\mathbf{w}_J\sqrt{P_J}\,n_J + n_k$',

  'eff_channel':
    r'$\mathbf{h}_{eff,k}[m] = \mathbf{G}^H[m]\,\Phi^H[m]\,'
    r'\mathbf{g}_{ru,k}[m] + \mathbf{g}_{bu,k}[m]$',

  'reflect_matrix':
    r'$\Phi[m] = \mathrm{diag}\!\left('
    r'e^{j\theta_1[m]},\ldots,e^{j\theta_M[m]}\right)$',

  'sinr':
    r'$\mathrm{SINR}_k[m] = \dfrac{P_k\,|\mathbf{h}_{eff,k}^H[m]\,'
    r'\mathbf{w}_k[m]|^2}{\sum_{i \neq k}P_i|\mathbf{h}_{eff,i}^H'
    r'\mathbf{w}_i|^2 + P_J|\mathbf{h}_{J,k}^H\mathbf{w}_J|^2 + \sigma^2}$',

  'rate_sum':
    r'$R_{sum} = \sum_{k=1}^{K}\frac{1}{M_{sc}}\sum_{m=1}^{M_{sc}}'
    r'\log_2\!\left(1+\mathrm{SINR}_k[m]\right)$',

  'protection':
    r'$\eta_{prot} = \frac{1}{K\,M_{sc}}\sum_k\sum_m'
    r'\mathbf{1}\!\left\{\mathrm{SINR}_k[m]\geq\gamma_{min}\right\}\times100\%$',

  'opt_problem':
    r'$\max_{\{P_k\},\,\Phi[m]} R_{sum}\;\;'
    r'\text{s.t.}\;\;\sum_k P_k\!\leq\!P_{max},\;\;'
    r'\mathrm{SINR}_k\!\geq\!\gamma_{min},\;\;|\theta_i|\!=\!1$',

  'path_loss':
    r'$PL(f,d)=\left(\dfrac{4\pi fd}{c}\right)^{\!2}\cdot e^{\,\kappa_{abs}(f)\cdot d}$',

  'channel_model':
    r'$\mathbf{G}[m]=\sqrt{\dfrac{NM}{L}}\sum_{\ell=1}^{L}'
    r'\alpha_\ell(f_m)\,\mathbf{a}_{RIS}(\theta_\ell^r)\,\mathbf{a}_{BS}^H(\theta_\ell^t)$',

  'beam_squint':
    r'$\eta(f_m)=\dfrac{1}{N}\,\left|\sum_{n=0}^{N-1}'
    r'\exp\!\left(j\pi n\cos\theta_0\cdot\dfrac{f_m-f_c}{f_c}\right)\right|^{\!2}$',

  'spdp_phase':
    r'$\theta_{q,i}[m]=\theta_{q,i}^{base}+2\pi f_m\tau_q$',

  'beamformer':
    r'$\mathbf{w}_k[m]=\mathbf{F}_{RF}\,\mathbf{f}_{BB,k}[m]$',

  'interf_cov':
    r'$\mathbf{R}_k=\sigma^2\mathbf{I}+\sum_{i\neq k}P_i\,'
    r'\tilde{\mathbf{h}}_i\,\tilde{\mathbf{h}}_i^H$',

  'mvdr':
    r'$\mathbf{f}_{BB,k}=\dfrac{\mathbf{R}_k^{-1}\tilde{\mathbf{h}}_k}'
    r'{\left\|\mathbf{R}_k^{-1}\tilde{\mathbf{h}}_k\right\|}$',

  'reward':
    r'$r=R_{sum}-0.5\,P_{frac}-3.0\sum_k'
    r'\mathbf{1}\!\left\{\mathrm{SINR}_k<\gamma_{min}\right\}$',

  'feature_pj':
    r'$f_{pj}=0.6\cdot\dfrac{\bar{P}_J-15}{25}+0.4\cdot\dfrac{\max P_J-15}{25}$',

  'feature_sinr':
    r'$f_{sinr}=0.5\cdot\dfrac{\overline{\mathrm{SINR}}+10}{40}'
    r'+0.5\cdot\dfrac{\min\mathrm{SINR}+10}{40}$',

  'fuzzy_mem':
    r'$\mu_i(x)=\max\!\left(0,\;1-\dfrac{|x-c_i|}{0.5}\right)$',

  'wolf_cond':
    r'$\text{if}\;\sum_a\pi(a)Q(s,a) > \sum_a\bar{\pi}(a)Q(s,a)'
    r'\;\Rightarrow\;\xi=\xi_{win}=0.01\;\;\text{else}\;\;\xi=\xi_{loss}=0.04$',

  'wolf_update':
    r'$\pi_\ell(a^*)\leftarrow\pi_\ell(a^*)+\xi\psi_\ell;\;\;'
    r'\pi_\ell(a)\leftarrow\pi_\ell(a)-\dfrac{\xi\psi_\ell}{|A|-1}\;\forall\, a\neq a^*$',

  'q_update':
    r'$Q_\ell\leftarrow Q_\ell+\alpha\psi_\ell\!\left['
    r'r+\gamma\max_{a^\prime}FQ(s^\prime,a^\prime)-Q_\ell\right]$',

  'fuzzy_q':
    r'$FQ(s,a)=\sum_\ell\psi_\ell\cdot Q_\ell(s,a)$',

  'predictability':
    r'$\eta=0.25\,\rho_{rep}+0.35\,\rho_{dom}+0.40\,(1-H_{norm})$',

  'jammer_power':
    r'$P_{J,k}[\mathrm{dBm}]=15+0.25\cdot\mathrm{clip}(\cdot)'
    r'+18\eta+\mathcal{N}(0,\,1.5^2)$',
}

def E(key): return EQ[key]  # shorthand

# ══════ SLIDE 1 – Title ══════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=LGRAY)
add_rect(sl, 0, 0, 13.33, 2.4, fill=NAVY)
add_rect(sl, 0, 2.4, 13.33, 0.06, fill=BLUE)
tb(sl, "Intelligent Reflecting Surface Assisted Anti-Jamming\nCommunications in THz Wideband Systems",
   0.5, 0.22, 12.3, 1.5, fs=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, "A Fuzzy WoLF-PHC Learning Approach with Hybrid Beamforming",
   0.5, 1.6, 12.3, 0.65, fs=17, color=RGBColor(0xBB,0xCC,0xEE),
   align=PP_ALIGN.CENTER, italic=True)
tb(sl, "Anuprabh Singh", 0.5, 2.8, 12.3, 0.45,
   fs=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
tb(sl, "Department of Electronics and Communication Engineering\n"
       "National Institute of Technology, Warangal",
   0.5, 3.22, 12.3, 0.7, fs=15, color=DGRAY, align=PP_ALIGN.CENTER)
add_rect(sl, 0.3, 4.1, 12.73, 0.025, fill=BLUE)
tb(sl, "Topics: System Model  •  THz Channel & Beam Squint  •  SPDP-RIS  •  "
       "Hybrid Beamforming  •  Fuzzy WoLF-PHC  •  Smart Jammer  •  Results",
   0.5, 4.25, 12.3, 0.6, fs=13, color=DGRAY, align=PP_ALIGN.CENTER)
tb(sl, "April 2026", 0.5, 6.9, 12.3, 0.4, fs=13,
   color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)

# ══════ SLIDE 2 – Motivation ══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "Motivation & Problem Statement",
    "Why this problem matters for 6G security")
buls(sl, ["Wireless channels are broadcast — any receiver can intercept or jam",
          "Smart jammers in 6G/THz are adaptive: they learn from the defender's actions",
          "THz bands (100 GHz) offer massive bandwidth but face unique challenges:",
          "  Severe molecular absorption & path loss at 100 GHz",
          "  Beam squint: 10 GHz bandwidth causes phased-array beams to mis-steer",
          "  Large arrays require expensive hybrid beamforming hardware"],
       0.35, 1.2, 7.6, 3.4, fs=15, title="The Threat Landscape", tc=NAVY)
add_rect(sl, 8.3, 1.2, 4.7, 3.4, fill=RGBColor(0xE8,0xF0,0xFB),
         line=BLUE, lw=Pt(1))
buls(sl, ["Narrowband (sub-6 GHz) only", "No beam squint compensation",
          "No hybrid beamforming", "Fixed or slow-adapting jammers",
          "No unpredictability guarantee"],
       8.45, 1.3, 4.4, 2.3, fs=14, title="Gaps in Prior Work", tc=RED)
add_rect(sl, 8.3, 4.0, 4.7, 0.85, fill=NAVY)
tb(sl, "This Paper: First THz IRS anti-jamming system\nwith Fuzzy WoLF-PHC + SPDP-RIS",
   8.4, 4.05, 4.5, 0.75, fs=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, "Three Core Challenges:", 0.35, 4.72, 8.0, 0.3, fs=13, bold=True, color=NAVY)
for i, (lbl, col) in enumerate([("Smart Adaptive\nJammer", RED),
                                  ("THz Beam\nSquint", BLUE),
                                  ("Joint Power +\nIRS Optimization", NAVY)]):
    x = 0.35 + i*2.65
    add_rect(sl, x, 4.85, 2.4, 0.95, fill=col)
    tb(sl, lbl, x+0.05, 4.9, 2.3, 0.85, fs=14, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)

# ══════ SLIDE 3 – System Architecture ═════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "System Architecture",
    "Wideband THz IRS-Assisted Multi-User Downlink")

params = [("Component","Parameter","Value"),
          ("Base Station","Antennas N","64 ULA"),
          ("Base Station","RF Chains N_RF","8"),
          ("IRS","Elements M","256 (16×16)"),
          ("IRS","Sub-arrays Q","64 SPDP"),
          ("OFDM","Subcarriers M_sc","64"),
          ("OFDM","Bandwidth B","10 GHz @ 100 GHz"),
          ("Users","Count K","4 single-antenna"),
          ("Jammer","Antennas N_J","2, adaptive"),
          ("System","P_max","40 dBm"),
          ("System","γ_min","5 dB"),
          ("System","Noise σ²","−82 dBm")]
cw = [2.2, 2.5, 1.8]
ch(sl,"Component",0.3,1.18,2.2); ch(sl,"Parameter",2.5,1.18,2.5); ch(sl,"Value",5.0,1.18,1.8)
for r, row in enumerate(params[1:]):
    tr(sl, row, 0.3, 1.56+r*0.38, cw, h=0.37, bg=LGRAY if r%2==0 else WHITE)

# Signal model equations — right panel
add_rect(sl, 7.1, 1.18, 5.9, 2.9, fill=RGBColor(0xF5,0xF8,0xFF), line=BLUE, lw=Pt(1))
tb(sl, "Received Signal", 7.2, 1.22, 5.7, 0.35, fs=13, bold=True, color=NAVY)
add_eq(sl, E('received_signal'), 7.12, 1.55, 5.86, 0.72, fs=13)
tb(sl, "Effective Channel & IRS Reflection", 7.2, 2.28, 5.7, 0.32, fs=13, bold=True, color=NAVY)
add_eq(sl, E('eff_channel'),     7.12, 2.60, 5.86, 0.52, fs=13)
add_eq(sl, E('reflect_matrix'),  7.12, 3.13, 5.86, 0.52, fs=13)

# SINR & objectives — bottom strip
add_rect(sl, 0.3, 5.62, 12.73, 1.6, fill=RGBColor(0xFF,0xF8,0xEC),
         line=RGBColor(0xFF,0x99,0x00), lw=Pt(1))
tb(sl, "SINR & Objectives", 0.45, 5.66, 5.0, 0.35, fs=13, bold=True,
   color=RGBColor(0xCC,0x66,0x00))
add_eq(sl, E('sinr'),        0.32, 5.98, 6.0, 0.62, fs=13, bg='#FFF8EC')
add_eq(sl, E('rate_sum'),    6.35, 5.82, 6.6, 0.55, fs=13, bg='#FFF8EC')
add_eq(sl, E('protection'),  6.35, 6.40, 6.6, 0.62, fs=13, bg='#FFF8EC')

# Optimization problem
add_rect(sl, 7.1, 4.08, 5.9, 1.42, fill=RGBColor(0xF5,0xEC,0xFF),
         line=RGBColor(0x88,0x44,0xBB), lw=Pt(1))
tb(sl, "Optimization Problem (NP-Hard)", 7.2, 4.12, 5.7, 0.35, fs=13, bold=True,
   color=RGBColor(0x66,0x22,0xAA))
add_eq(sl, E('opt_problem'), 7.12, 4.46, 5.86, 1.0, fs=12, bg='#F5ECFF')

# Geometry
add_rect(sl, 0.3, 4.08, 6.5, 1.42, fill=RGBColor(0xF0,0xF8,0xF0), line=GREEN, lw=Pt(1))
tb(sl, "Network Geometry", 0.45, 4.12, 6.2, 0.35, fs=13, bold=True, color=GREEN)
tb(sl, "BS @ (0,0)  →  IRS @ (50,75) m  →  Users 30–120 m\n"
       "Jammer: random walk [40–100]×[0–80] m\n"
       "Rician K: K_BR=8 dB, K_BU=3 dB, K_RU=6 dB",
   0.45, 4.5, 6.2, 0.95, fs=12.5, color=DGRAY)

# ══════ SLIDE 4 – THz Channel & Beam Squint ═══════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "THz Channel Model & Beam Squint Problem",
    "Why 10 GHz bandwidth at 100 GHz demands a special IRS architecture")

tb(sl, "THz Channel Model:", 0.35, 1.18, 8.5, 0.38, fs=15, bold=True, color=NAVY)
add_eq(sl, E('path_loss'),     0.35, 1.55, 5.8, 0.58, fs=17)
add_eq(sl, E('channel_model'), 0.35, 2.13, 8.5, 0.62, fs=16)
tb(sl, "Beam Squint — Normalized Array Gain:", 0.35, 2.82, 8.5, 0.38,
   fs=15, bold=True, color=NAVY)
add_eq(sl, E('beam_squint'), 0.35, 3.2, 8.5, 0.78, fs=16)
tb(sl, "At β = B/f_c = 10%:  only the center few subcarriers get near-full gain. "
       "~61 of 64 subcarriers lose almost all array gain.",
   0.35, 4.02, 8.5, 0.5, fs=13, color=DGRAY)

tb(sl, "Beam Squint Severity (N=64, θ₀=30°):", 0.35, 4.58, 8.5, 0.38,
   fs=14, bold=True, color=NAVY)
bww = [1.8, 1.6, 1.6, 1.8]
ch(sl,"Bandwidth",0.3,4.98,1.8,fs=12); ch(sl,"β",2.1,4.98,1.6,fs=12)
ch(sl,"Edge Gain",3.7,4.98,1.6,fs=12); ch(sl,"Verdict",5.3,4.98,1.8,fs=12)
for r, row in enumerate([("0.1 GHz","0.1%","~100%","✓ Perfect"),
                          ("2 GHz","2%","~40%","△ Manageable"),
                          ("10 GHz","10%","~1%","✗ Catastrophic")]):
    bg = RGBColor(0xFF,0xEE,0xEE) if r==2 else (LGRAY if r%2==0 else WHITE)
    tr(sl, row, 0.3, 5.36+r*0.38, bww, h=0.37, bg=bg,
       fg=RED if r==2 else DGRAY, bold=(r==2))

add_rect(sl, 9.1, 1.18, 4.0, 5.2, fill=RGBColor(0xF2,0xF8,0xFF), line=BLUE, lw=Pt(1))
tb(sl, "Figure 1 (Paper)", 9.2, 1.22, 3.8, 0.35, fs=14, bold=True,
   color=NAVY, align=PP_ALIGN.CENTER)
tb(sl, "Plot: η(f_m) vs. subcarrier index m\n\n"
       "B=0.1 GHz → flat line at 1.0\n"
       "B=2 GHz   → gentle arch\n"
       "B=10 GHz  → spike at centre only\n\n"
       "Motivates SPDP-RIS with\nTrue-Time Delay elements.",
   9.2, 1.62, 3.8, 2.55, fs=13, color=DGRAY, align=PP_ALIGN.CENTER)
add_eq(sl, E('beam_squint'), 9.12, 4.32, 3.86, 1.7, fs=13,
       bg='#F2F8FF', tc='#1A3A6C')
tb(sl, "→ Motivates SPDP-RIS with TTD elements",
   0.35, 6.52, 8.5, 0.38, fs=14, bold=True, color=NAVY)

# ══════ SLIDE 5 – SPDP-RIS & Hybrid BF ═══════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "SPDP-RIS & Hybrid Beamforming at BS",
    "Hardware solutions for wideband THz beam squint compensation")

add_rect(sl, 0.3, 1.18, 6.1, 5.9, fill=RGBColor(0xF2,0xF8,0xFF), line=BLUE, lw=Pt(1))
tb(sl, "SPDP-RIS Architecture", 0.45, 1.22, 5.8, 0.38, fs=16, bold=True, color=NAVY)
buls(sl, ["256 IRS elements → Q=64 sub-arrays (4 elements each)",
          "Each sub-array: 1 TTD element + per-element phase shifters",
          "TTD gives frequency-proportional phase: no squint",
          "Cost: 64 TTDs only (vs 256 for full per-element TTD)",
          "Result: all 64 subcarriers get near-full array gain"],
       0.45, 1.65, 5.8, 1.6, fs=13)
tb(sl, "Phase per element:", 0.45, 3.32, 5.8, 0.32, fs=13, bold=True, color=NAVY)
add_eq(sl, E('spdp_phase'), 0.35, 3.62, 6.0, 0.62, fs=18)
tb(sl, "θ_base: centre-frequency steering phase\n"
       "τ_q: sub-array TTD delay (compensates squint)",
   0.45, 4.28, 5.8, 0.5, fs=12.5, color=DGRAY)
add_rect(sl, 0.45, 4.85, 5.7, 0.95, fill=NAVY)
tb(sl, "SPDP Q=64 closely matches Classical (no-squint)\nacross all 64 subcarriers — confirmed in Figure 1.",
   0.55, 4.9, 5.5, 0.85, fs=12.5, color=WHITE)

add_rect(sl, 6.7, 1.18, 6.3, 5.9, fill=RGBColor(0xF5,0xF0,0xFF),
         line=RGBColor(0x88,0x44,0xBB), lw=Pt(1))
tb(sl, "Hybrid Beamforming at BS", 6.85, 1.22, 6.0, 0.38, fs=16, bold=True,
   color=RGBColor(0x66,0x22,0xAA))
buls(sl, ["64 antennas, 8 RF chains (sub-connected architecture)",
          "Each RF chain drives 8 antennas independently"],
       6.85, 1.65, 6.0, 0.7, fs=13)
tb(sl, "Beamformer decomposition:", 6.85, 2.42, 6.0, 0.32, fs=13, bold=True, color=NAVY)
add_eq(sl, E('beamformer'), 6.72, 2.72, 6.26, 0.58, fs=18, bg='#F5F0FF')
tb(sl, "F_RF ∈ ℂ^{64×8}: analog precoder (phase shifters)\n"
       "f_{BB,k}[m] ∈ ℂ^8: digital precoder per subcarrier",
   6.85, 3.35, 6.0, 0.55, fs=12.5, color=DGRAY)
tb(sl, "Interference covariance:", 6.85, 3.98, 6.0, 0.32, fs=13, bold=True, color=NAVY)
add_eq(sl, E('interf_cov'), 6.72, 4.28, 6.26, 0.62, fs=18, bg='#F5F0FF')
tb(sl, "MVDR digital precoder:", 6.85, 4.98, 6.0, 0.32, fs=13, bold=True, color=NAVY)
add_eq(sl, E('mvdr'), 6.72, 5.28, 6.26, 0.72, fs=18, bg='#F5F0FF')
add_rect(sl, 6.85, 6.1, 5.7, 0.72, fill=RGBColor(0x66,0x22,0xAA))
tb(sl, "8× hardware cost reduction vs full-digital at 100 GHz\n"
       "while retaining near-optimal spatial multiplexing.",
   6.95, 6.15, 5.5, 0.62, fs=12.5, color=WHITE)

# ══════ SLIDE 6 – RL Framework ═══════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "Reinforcement Learning Framework",
    "Hybrid RL + Alternating Optimization decomposition")
tb(sl, "Why RL?  —  Three Reasons:", 0.35, 1.18, 12.6, 0.38, fs=16, bold=True, color=NAVY)
for i, (reason, detail) in enumerate([
        ("Unknown Jammer", "Adversarial non-stationary strategy — no closed-form model"),
        ("Non-Convex Problem","Joint P_k (discrete) + Φ[m] (continuous) is NP-hard"),
        ("Adaptive Environment","Channels & jammer position change every episode")]):
    x = 0.35 + i*4.3
    add_rect(sl, x, 1.62, 4.1, 0.9, fill=NAVY)
    tb(sl, reason, x+0.1, 1.65, 3.9, 0.35, fs=14, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, detail, x+0.1, 2.0, 3.9, 0.48, fs=12,
       color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)

tb(sl, "Hybrid RL + AO Pipeline:", 0.35, 2.72, 12.6, 0.38, fs=15, bold=True, color=NAVY)
for i, (title, body) in enumerate([
        ("① State\nObservation","Read P_J, channel,\nSINR → 3 features\n∈ [0,1]³"),
        ("② RL Action","Pick 1 of 42:\n7 power fracs ×\n6 alloc modes"),
        ("③ AO: IRS\nPhases","Given {P_k}, solve\n256 phases via\nAO (6 iter)"),
        ("④ MVDR\nBF","Compute {w_k}\nvia R_k^{-1}"),
        ("⑤ Reward\n& Update","r = R_sum − penalty\nUpdate Q & π")]):
    x = 0.3 + i*2.55
    add_rect(sl, x, 3.18, 2.4, 1.5, fill=LGRAY, line=BLUE, lw=Pt(1))
    tb(sl, title, x+0.05, 3.22, 2.3, 0.45, fs=13, bold=True,
       color=NAVY, align=PP_ALIGN.CENTER)
    tb(sl, body, x+0.05, 3.65, 2.3, 0.95, fs=11.5,
       color=DGRAY, align=PP_ALIGN.CENTER)
    if i < 4:
        tb(sl, "→", x+2.28, 3.72, 0.3, 0.4, fs=22, bold=True,
           color=BLUE, align=PP_ALIGN.CENTER)

add_rect(sl, 0.3, 4.85, 12.73, 1.05, fill=LGRAY, line=NAVY, lw=Pt(0.5))
tb(sl, "Reward Function:", 0.45, 4.9, 3.5, 0.35, fs=14, bold=True, color=NAVY)
add_eq(sl, E('reward'), 0.35, 5.22, 12.6, 0.62, fs=20, bg='#F2F4F7')

tb(sl, "Action Space — 42 actions  (7 power fractions × 6 allocation modes):",
   0.35, 6.05, 12.6, 0.35, fs=14, bold=True, color=NAVY)
for i, m in enumerate(["equal  P_k=P/K", "channel-prop.  P_k∝‖h‖²",
                        "inverse-ch.  P_k∝1/‖h‖²",
                        "SINR-deficit  P_k∝(γ_min−SINR_k)⁺",
                        "water-filling", "max-min fair"]):
    x = 0.35 + (i%3)*4.1; y = 6.44 + (i//3)*0.42
    add_rect(sl, x, y, 3.9, 0.38, fill=BLUE)
    tb(sl, m, x+0.08, y+0.05, 3.75, 0.3, fs=12, color=WHITE)

# ══════ SLIDE 7 – State & Fuzzy Aggregation ═══════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "State Representation & Fuzzy Aggregation",
    "Three normalized features → 27-component fuzzy state vector")
tb(sl, "Three Normalized Features  f = [f_pj,  f_ch,  f_sinr] ∈ [0,1]³",
   0.35, 1.18, 12.6, 0.4, fs=16, bold=True, color=NAVY)

for i, (name, col, key, note) in enumerate([
        ("Jammer Pressure  f_pj", BLUE, 'feature_pj',
         "0 = weak jammer          1 = full-power  (60/40 sustained + burst)"),
        ("SINR Health  f_sinr", RED, 'feature_sinr',
         "0 = all users failing    1 = all comfortable  (50/50 avg + worst-case)")]):
    y = 1.65 + i*2.1
    add_rect(sl, 0.3, y, 0.15, 1.7, fill=col)
    tb(sl, name, 0.55, y+0.02, 4.5, 0.38, fs=14, bold=True, color=col)
    add_eq(sl, E(key), 0.55, y+0.44, 8.7, 0.78, fs=20)
    tb(sl, note, 0.55, y+1.26, 8.7, 0.4, fs=12.5, color=DGRAY)

add_rect(sl, 0.3, 3.75, 0.15, 1.7, fill=GREEN)
tb(sl, "Channel Quality  f_ch", 0.55, 3.77, 4.5, 0.38, fs=14, bold=True, color=GREEN)
tb(sl, "f_ch = 0.75·(mean_ch_dB + 100)/60  +  0.25·(1 − std_spread/20)\n"
       "Range: [−100 dB, −40 dB] → [0,1]\n"
       "0 = deep fade     1 = strong uniform channels across all users",
   0.55, 4.18, 8.7, 1.2, fs=12.5, color=DGRAY)

add_rect(sl, 9.5, 1.65, 3.5, 5.6, fill=RGBColor(0xF2,0xF8,0xFF), line=BLUE, lw=Pt(1))
tb(sl, "Fuzzy Aggregation", 9.6, 1.7, 3.3, 0.38, fs=14, bold=True,
   color=NAVY, align=PP_ALIGN.CENTER)
tb(sl, "3 triangular kernels per feature\nCenters: {0.0, 0.5, 1.0}",
   9.6, 2.15, 3.3, 0.55, fs=13, color=DGRAY, align=PP_ALIGN.CENTER)
add_eq(sl, E('fuzzy_mem'), 9.52, 2.75, 3.46, 0.88, fs=15, bg='#F2F8FF')
tb(sl, "Joint 3D product:\n"
       "ψ_ℓ = μ_i(f_pj)·μ_j(f_ch)·μ_k(f_sinr)\n"
       "ℓ = 9i + 3j + k\n\n"
       "27 fuzzy components  (3³)\nΣ_ℓ ψ_ℓ = 1.0\n\n"
       "512 discrete state IDs  (8³)\n\n"
       "Key benefit:\nNearby states share Q-values\n→ 14× lower variance",
   9.6, 3.68, 3.3, 3.52, fs=12.5, color=DGRAY, align=PP_ALIGN.CENTER)

# ══════ SLIDE 8 – WoLF-PHC ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "Fuzzy WoLF-PHC Algorithm",
    "Win or Learn Fast — Policy Hill Climbing with fuzzy state aggregation")

add_rect(sl, 0.3, 1.18, 5.9, 2.5, fill=RGBColor(0xFF,0xF3,0xF3), line=RED, lw=Pt(1))
tb(sl, "Why Not Standard Q-Learning?", 0.45, 1.22, 5.6, 0.38, fs=14, bold=True, color=RED)
tb(sl, "Q-learning → deterministic policy → η≈0.85 → +15.3 dB jammer boost\n"
       "WoLF-PHC → mixed strategy  → η≈0.13 → +2.3 dB jammer boost\n\n"
       "12 dB gap = 16× more jamming power against Q-learning\n\n"
       "Jammer monitors last 25 actions to compute η and exploits predictability.",
   0.45, 1.65, 5.6, 1.98, fs=12.5, color=DGRAY)

add_rect(sl, 6.5, 1.18, 6.5, 2.5, fill=RGBColor(0xF2,0xF8,0xFF), line=BLUE, lw=Pt(1))
tb(sl, "WoLF Mechanism — Adaptive Learning Rate", 6.65, 1.22, 6.2, 0.38,
   fs=14, bold=True, color=NAVY)
tb(sl, "Maintains: Q_ℓ(s,a),  π_ℓ(a),  π̄_ℓ(a)  per fuzzy component",
   6.65, 1.62, 6.2, 0.32, fs=12.5, color=DGRAY)
add_eq(sl, E('wolf_cond'),   6.52, 1.96, 6.46, 0.7,  fs=13, bg='#F2F8FF')
tb(sl, "Policy update:", 6.65, 2.72, 3.5, 0.32, fs=13, bold=True, color=NAVY)
add_eq(sl, E('wolf_update'), 6.52, 3.02, 6.46, 0.62, fs=14, bg='#F2F8FF')

tb(sl, "Fuzzy Q-Value Update (Bellman equation):", 0.35, 3.82, 12.6, 0.38,
   fs=15, bold=True, color=NAVY)
add_rect(sl, 0.3, 4.2, 12.73, 0.82, fill=LGRAY, line=NAVY, lw=Pt(0.5))
add_eq(sl, E('q_update'), 0.35, 4.22, 12.6, 0.78, fs=20, bg='#F2F4F7')

tb(sl, "Fuzzy Q-value:", 0.35, 5.12, 3.0, 0.32, fs=13, bold=True, color=NAVY)
add_eq(sl, E('fuzzy_q'), 0.35, 5.45, 5.8, 0.55, fs=20)

tb(sl, "Key Enhancements:", 0.35, 6.08, 12.6, 0.32, fs=15, bold=True, color=NAVY)
for i, (title, body) in enumerate([
        ("Adaptive α","α·(1+3/(1+0.1·N_visits))\n4× faster for new pairs"),
        ("Experience Replay","256 buffer, 4 replays/step"),
        ("Boltzmann Eval","Softmax with adaptive temp"),
        ("3% Uniform Floor","Every action ≥ 0.07% prob")]):
    x = 0.3 + i*3.25
    add_rect(sl, x, 6.42, 3.1, 0.9, fill=LGRAY, line=BLUE, lw=Pt(1))
    tb(sl, title, x+0.08, 6.45, 2.95, 0.32, fs=13, bold=True,
       color=NAVY, align=PP_ALIGN.CENTER)
    tb(sl, body, x+0.08, 6.78, 2.95, 0.5, fs=12,
       color=DGRAY, align=PP_ALIGN.CENTER)

# ══════ SLIDE 9 – Smart Jammer ════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "Smart Jammer Model",
    "Predictability-exploiting adaptive adversary")
tb(sl, "Predictability Score:", 0.35, 1.18, 8.0, 0.38, fs=15, bold=True, color=NAVY)
add_eq(sl, E('predictability'), 0.35, 1.56, 8.5, 0.62, fs=20)
tb(sl, "ρ_rep = repetition rate of consecutive actions\n"
       "ρ_dom = fraction of actions = most common action\n"
       "H_norm = normalized entropy (1 = fully random, 0 = deterministic)",
   0.45, 2.22, 8.5, 0.65, fs=12.5, color=DGRAY)
tb(sl, "Jammer Power Adaptation:", 0.35, 3.0, 8.0, 0.38, fs=15, bold=True, color=NAVY)
add_rect(sl, 0.3, 3.38, 12.73, 0.88, fill=RGBColor(0xFF,0xF0,0xF0),
         line=RED, lw=Pt(0.5))
add_eq(sl, E('jammer_power'), 0.35, 3.38, 12.6, 0.88, fs=18, bg='#FFF0F0', tc='#222222')
tb(sl, "EMA smoothing:  P_J_smooth = 0.7·P_J_prev + 0.3·P_J_target",
   0.45, 4.3, 10.0, 0.35, fs=13, color=DGRAY)

tb(sl, "Impact of Predictability on Jammer Strength:", 0.35, 4.78, 12.6, 0.38,
   fs=15, bold=True, color=NAVY)
cw2 = [2.4, 1.4, 1.5, 2.3, 2.4]
ch(sl,"Agent",0.3,5.18,2.4,fs=12); ch(sl,"η",2.7,5.18,1.4,fs=12)
ch(sl,"Exploit Boost",4.1,5.18,1.5,fs=12)
ch(sl,"Eff. P_J",5.6,5.18,2.3,fs=12); ch(sl,"Precoder",7.9,5.18,2.4,fs=12)
for r, row in enumerate([
        ("Fuzzy WoLF-PHC","≈0.13","+2.3 dB","≈17.3 dBm","Random"),
        ("DQN","≈0.60","+10.8 dB","≈25.8 dBm","Partial"),
        ("Q-learning","≈0.85","+15.3 dB","≈30.3 dBm","Mostly aligned"),
        ("AO Baseline","1.00","+18.0 dB","≈33.0 dBm","Fully aligned")]):
    bg = RGBColor(0xE8,0xF8,0xE8) if r==0 else (RGBColor(0xFF,0xEE,0xEE) if r==3 else (LGRAY if r%2==0 else WHITE))
    fg = GREEN if r==0 else (RED if r==3 else DGRAY)
    tr(sl, row, 0.3, 5.56+r*0.38, cw2, h=0.37, bg=bg, fg=fg, bold=(r==0 or r==3))
add_rect(sl, 0.3, 7.05, 12.73, 0.33, fill=NAVY)
tb(sl, "16 dB gap = 40× more jamming power vs AO Baseline — this mechanism explains most performance differences.",
   0.45, 7.08, 12.3, 0.28, fs=13, bold=True, color=WHITE)

# ══════ SLIDE 10 – Novelty ════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "Novelty & Contributions",
    "What makes this work unique compared to prior art")
tb(sl, "Prior Work (Yang et al. 2021): IRS anti-jamming with Fuzzy WoLF-PHC — "
       "narrowband sub-6 GHz, no beam squint, no hybrid beamforming, no THz.",
   0.35, 1.18, 12.6, 0.55, fs=13, color=RGBColor(0x66,0x44,0x00))
for i, (title, color, body) in enumerate([
        ("①  First THz IRS Anti-Jamming System", NAVY,
         "Extends IRS anti-jamming to 100 GHz / 10 GHz bandwidth — first to address THz adversarial "
         "channel with beam squint, molecular absorption, and wideband OFDM simultaneously."),
        ("②  SPDP-RIS with True-Time Delay", BLUE,
         "Proposes SPDP architecture with TTD elements to compensate beam squint. "
         "Reduces hardware from 256 to 64 TTDs while maintaining near-ideal wideband gain."),
        ("③  Sub-Connected Hybrid Beamforming", GREEN,
         "Integrates hybrid analog-digital BF at 64-antenna BS with MVDR digital precoding "
         "per subcarrier, enabling practical 100 GHz hardware with only 8 RF chains."),
        ("④  Unpredictability as Security Property", RGBColor(0x88,0x22,0xAA),
         "Formally demonstrates WoLF-PHC's mixed strategy keeps η≈0.13 regardless of jammer power, "
         "limiting exploitation to +2.3 dB vs +15–18 dB for all deterministic baselines.")]):
    y = 1.88 + i*1.25
    add_rect(sl, 0.3, y, 0.25, 1.05, fill=color)
    tb(sl, title, 0.65, y+0.02, 5.4, 0.38, fs=14, bold=True, color=color)
    tb(sl, body, 0.65, y+0.45, 12.35, 0.55, fs=13, color=DGRAY)

# ══════ SLIDE 11 – Results: Overall ══════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "Results: Overall Performance — Table II & Figure 3",
    "Mean ± std over 3 seeds × 50 eval episodes × 20 steps")
cw3 = [3.5, 2.0, 1.0, 2.0, 1.0, 1.5]; x0 = 0.3
ch(sl,"Method",x0,1.2,3.5); ch(sl,"Rate (b/s/Hz)",x0+3.5,1.2,2.0)
ch(sl,"Std",x0+5.5,1.2,1.0); ch(sl,"SINR Prot.",x0+6.5,1.2,2.0)
ch(sl,"Std",x0+8.5,1.2,1.0); ch(sl,"Jammer η",x0+9.5,1.2,1.5)
for r, row in enumerate([
        ("Fuzzy WoLF-PHC (Proposed)","11.46","±0.18","78.6%","±1.9%","≈0.13"),
        ("AO Baseline [Wu'19]","8.01","±0.53","44.8%","±4.8%","1.00"),
        ("DQN [Mnih'15]","7.68","±1.14","41.9%","±11.4%","≈0.60"),
        ("Fast Q-Learning [Yang'21]","7.21","±1.97","39.2%","±19.1%","≈0.80"),
        ("Classical Q-Learning","7.13","±2.52","38.8%","±20.9%","≈0.85"),
        ("No IRS (Lower Bound)","5.48","—","18.8%","—","—")]):
    if r==0: bg=RGBColor(0xD4,0xED,0xFF); fg=NAVY; bld=True
    elif r==5: bg=RGBColor(0xF5,0xF5,0xF5); fg=DGRAY; bld=False
    else: bg=LGRAY if r%2==0 else WHITE; fg=DGRAY; bld=False
    tr(sl, row, x0, 1.58+r*0.4, cw3, h=0.39, bg=bg, fg=fg, bold=bld)
tb(sl, "Key Observations:", 0.35, 4.45, 12.6, 0.38, fs=15, bold=True, color=NAVY)
for i, (title, body, color) in enumerate([
        ("43% higher rate","11.46 vs 8.01 (AO Baseline)",BLUE),
        ("75% higher protection","78.6% vs 44.8%",GREEN),
        ("14× lower variance","±0.18 vs ±2.52",NAVY),
        ("2× IRS gain","5.48 (No-IRS) → 11.46",RGBColor(0x88,0x22,0xAA))]):
    x = 0.3+(i%2)*6.5; y = 4.88+(i//2)*1.1
    add_rect(sl, x, y, 6.3, 1.0, fill=LGRAY, line=color, lw=Pt(1.5))
    tb(sl, title, x+0.12, y+0.05, 6.0, 0.38, fs=14, bold=True, color=color)
    tb(sl, body, x+0.12, y+0.48, 6.0, 0.45, fs=13, color=DGRAY)
tb(sl, "Root cause: 16 dB more jammer power vs AO (η=1.0) → 40× jamming intensity",
   0.35, 7.08, 12.6, 0.32, fs=13, bold=True, color=RED, align=PP_ALIGN.CENTER)

# ══════ SLIDE 12 – Convergence & Per-Seed ════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "Results: Convergence (Fig. 2) & Per-Seed Consistency (Fig. 10)",
    "Learning stability and deployment reliability across random seeds")
add_rect(sl, 0.3, 1.18, 6.1, 5.6, fill=LGRAY, line=BLUE, lw=Pt(1))
tb(sl, "Figure 2 — Convergence Curves", 0.45, 1.22, 5.8, 0.38, fs=15, bold=True, color=NAVY)
tb(sl, "X: Training episode (0–1200+)\nY: Average reward per episode\n\n"
       "WoLF-PHC (solid blue): gradually rises as mixed\n"
       "policy keeps η low throughout training.\n\n"
       "Q-learning plateaus then drops as ε→0 and\nthe jammer adapts to the deterministic policy.\n\n"
       "Shaded bands = min–max across 3 seeds:\n"
       "  WoLF-PHC: narrow → consistent\n"
       "  Q-learning: wide → unreliable\n\n"
       "WoLF-PHC trains 3× longer (3600 episodes)\nfor policy convergence.",
   0.45, 1.65, 5.8, 5.05, fs=13, color=DGRAY)
add_rect(sl, 6.7, 1.18, 6.3, 5.6, fill=LGRAY, line=BLUE, lw=Pt(1))
tb(sl, "Figure 10 — Per-Seed Results", 6.85, 1.22, 6.0, 0.38, fs=15, bold=True, color=NAVY)
sw3 = [2.5, 1.9, 1.7]
ch(sl,"Method",6.7,1.65,2.5,fs=12); ch(sl,"Rate Range",9.2,1.65,1.9,fs=12); ch(sl,"Prot. Range",11.1,1.65,1.7,fs=12)
for r, row in enumerate([("WoLF-PHC","11.3–11.6","76%–81%"),
                           ("Q-learning","3.5–9.3","10%–55%"),
                           ("Fast Q-learning","4.5–9.3","13%–61%"),
                           ("DQN","6.0–8.5","20%–55%"),
                           ("AO Baseline","7.5–8.5","40%–50%")]):
    bg = RGBColor(0xD4,0xED,0xFF) if r==0 else (LGRAY if r%2==0 else WHITE)
    tr(sl, row, 6.7, 2.03+r*0.38, sw3, h=0.37, bg=bg, fg=NAVY if r==0 else DGRAY, bold=(r==0))
tb(sl, "WoLF-PHC: rate variance 0.3 b/s/Hz\n"
       "Q-learning: rate variance 5.8 b/s/Hz\n→ 19× more consistent across seeds\n\n"
       "Deploying 'bad seed' Q-learning:\n3.5 b/s/Hz + 10% protection = unusable.\n\n"
       "WoLF-PHC guarantees ≥11.3 b/s/Hz\nand ≥76% on every deployment.",
   6.85, 2.85, 5.8, 2.2, fs=13, color=DGRAY)
tb(sl, "Why: fuzzy memberships share Q-values across nearby states → robust to any channel realization",
   0.35, 6.9, 12.6, 0.32, fs=14, bold=True, color=NAVY)

# ══════ SLIDE 13 – Power & IRS Sweeps ═════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "Results: Transmit Power Sweep (Fig. 4) & IRS Size Sweep (Fig. 5)",
    "Effect of P_max ∈ {25,32,40} dBm and N_RIS ∈ {16,64,144,256}")
tb(sl, "Figure 4 — Transmit Power Sweep:", 0.35, 1.18, 6.3, 0.38, fs=15, bold=True, color=NAVY)
pww = [1.3, 1.4, 1.2, 1.4, 1.5]
ch(sl,"P_max",0.3,1.6,1.3,fs=12); ch(sl,"WoLF R",1.6,1.6,1.4,fs=12)
ch(sl,"Q R",3.0,1.6,1.2,fs=12); ch(sl,"WoLF P",4.2,1.6,1.4,fs=12); ch(sl,"Q P",5.6,1.6,1.5,fs=12)
for r, row in enumerate([("25 dBm","1.04","0.44","0.0%","0.0%"),
                           ("32 dBm","3.49","3.30","1.5%","0.6%"),
                           ("40 dBm ✓","9.14","8.12","61.6%","50.6%")]):
    tr(sl, row, 0.3, 1.98+r*0.38, pww, h=0.37,
       bg=RGBColor(0xD4,0xED,0xFF) if r==2 else (LGRAY if r%2==0 else WHITE), bold=(r==2))
add_rect(sl, 0.3, 3.14, 7.0, 1.48, fill=RGBColor(0xF2,0xF8,0xFF), line=BLUE, lw=Pt(1))
tb(sl, "• 25 dBm: everyone fails — noise+jamming overpowers all links\n"
       "• 32 dBm: marginal — methods start separating\n"
       "• 40 dBm: full operation — WoLF-PHC 12 dB lower exploitation decisive\n"
       "• No-IRS (~5.5): IRS provides substantial gain at sufficient power",
   0.45, 3.18, 6.7, 1.38, fs=12.5, color=DGRAY)
tb(sl, "Figure 5 — IRS Size Sweep:", 7.4, 1.18, 5.6, 0.38, fs=15, bold=True, color=NAVY)
rw = [0.85, 1.3, 1.3, 1.1, 1.0]
ch(sl,"N_RIS",7.4,1.6,0.85,fs=12); ch(sl,"WoLF R",8.25,1.6,1.3,fs=12)
ch(sl,"WoLF P",9.55,1.6,1.3,fs=12); ch(sl,"Q R",10.85,1.6,1.1,fs=12); ch(sl,"AO R",11.95,1.6,1.0,fs=12)
for r, row in enumerate([("16","8.42","52.7%","8.05","5.48"),
                           ("64","8.47","53.3%","8.21","5.53"),
                           ("144","8.62","55.4%","8.28","5.69"),
                           ("256 ✓","9.00","59.6%","8.10","5.96")]):
    tr(sl, row, 7.4, 1.98+r*0.38, rw, h=0.37,
       bg=RGBColor(0xD4,0xED,0xFF) if r==3 else (LGRAY if r%2==0 else WHITE), bold=(r==3))
add_rect(sl, 7.4, 3.14, 5.6, 1.48, fill=RGBColor(0xF2,0xF8,0xFF), line=BLUE, lw=Pt(1))
tb(sl, "• 16 elements: 8.42 vs 5.48 No-IRS — 54% gain from tiny panel\n"
       "• 16→256: modest +0.58 (diminishing returns per element)\n"
       "• AO stays low: η=1.0 negates any IRS DoF gain\n"
       "• WoLF-PHC benefits from more elements AND jammer resistance",
   7.5, 3.18, 5.3, 1.38, fs=12.5, color=DGRAY)
tb(sl, "Combined Insights:", 0.35, 4.75, 12.6, 0.38, fs=15, bold=True, color=NAVY)
for i, t in enumerate([
        "Power threshold: below ~32 dBm jamming overwhelms all — sufficient TX power is a prerequisite",
        "IRS confirmed beneficial even at N=16 elements — hardware deployment feasibility demonstrated",
        "WoLF-PHC leads across ALL power and IRS size combinations — robustness is not parameter-dependent",
        "AO Baseline fails to scale with IRS size because η=1.0 negates additional beamforming DoF"]):
    tb(sl, f"▸  {t}", 0.35, 5.18+i*0.52, 12.6, 0.48, fs=13, color=DGRAY)

# ══════ SLIDE 14 – SINR & Jammer Power Sweeps ════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "Results: SINR Target Sweep (Fig. 6) & Jammer Power Sweep (Fig. 7)",
    "Effect of γ_min and P_J — the most critical results")
tb(sl, "Figure 6 — SINR Target Sweep:", 0.35, 1.18, 6.3, 0.38, fs=15, bold=True, color=NAVY)
sw4 = [0.9, 1.3, 1.4, 1.5, 1.4]
ch(sl,"γ_min",0.3,1.6,0.9,fs=12); ch(sl,"WoLF R",1.2,1.6,1.3,fs=12)
ch(sl,"WoLF P",2.5,1.6,1.4,fs=12); ch(sl,"Fast-Q R",3.9,1.6,1.5,fs=12); ch(sl,"Fast-Q P",5.4,1.6,1.4,fs=12)
for r, row in enumerate([("3 dB","9.13","78.8%","7.69","63.5%"),
                           ("10 dB","9.01","7.9%","4.55","0.8%"),
                           ("20 dB","8.95","~0%","4.27","0.0%")]):
    tr(sl, row, 0.3, 1.98+r*0.38, sw4, h=0.37,
       bg=RGBColor(0xD4,0xED,0xFF) if r==0 else (LGRAY if r%2==0 else WHITE))
add_rect(sl, 0.3, 3.15, 6.4, 1.88, fill=RGBColor(0xF2,0xF8,0xFF), line=BLUE, lw=Pt(1))
tb(sl, "• WoLF-PHC rate nearly FLAT (9.13→8.95): no throughput–protection tradeoff\n"
       "• Protection collapses at 10 dB: THz noise+jammer makes high SINR across K·M_sc targets hard\n"
       "• Fast-Q collapses 7.69→4.27 (−44%): stricter threshold destabilises greedy Q-table\n"
       "• AO rate constant (5.97): ignores threshold, determinism → always η=1.0",
   0.45, 3.2, 6.2, 1.78, fs=12.5, color=DGRAY)
tb(sl, "Figure 7 — Jammer Power Sweep  (THE KEY RESULT):", 7.0, 1.18, 6.0, 0.38,
   fs=15, bold=True, color=RED)
jw = [1.0, 1.3, 1.2, 1.4, 1.2]
ch(sl,"P_J max",7.0,1.6,1.0,fs=12); ch(sl,"WoLF R",8.0,1.6,1.3,fs=12)
ch(sl,"Q R",9.3,1.6,1.2,fs=12); ch(sl,"Fast-Q R",10.5,1.6,1.4,fs=12); ch(sl,"DQN R",11.9,1.6,1.2,fs=12)
for r, row in enumerate([("5 dBm","9.26","10.72","11.23","11.17"),
                           ("12 dBm","9.02","9.81","9.33","≈9.5"),
                           ("20 dBm","9.14 ✓","8.12 ↓","4.63 ✗","7.25 ↓")]):
    bg = [LGRAY, WHITE, RGBColor(0xD4,0xED,0xFF)][r]
    tr(sl, row, 7.0, 1.98+r*0.38, jw, h=0.37, bg=bg,
       fg=GREEN if r==2 else DGRAY, bold=(r==2))
add_rect(sl, 7.0, 3.15, 6.1, 1.88, fill=RGBColor(0xFF,0xF0,0xF0), line=RED, lw=Pt(1))
tb(sl, "① 5 dBm: WoLF-PHC starts LOWER — greedy is optimal for weak jammer,\n"
       "   WoLF-PHC's 15% stochastic component wastes some actions.\n\n"
       "② 12 dBm: separation begins. Fast-Q drops as jammer learns\n"
       "   pattern (η→0.8 → +14.4 dB boost).\n\n"
       "③ 20 dBm: THE DECISIVE RESULT:\n"
       "   Fast-Q: 4.63 (−59%!)  CATASTROPHIC COLLAPSE\n"
       "   WoLF-PHC: 9.14 (−1.3%)  ≈ FLAT LINE",
   7.1, 3.2, 5.8, 1.78, fs=12.5, color=DGRAY)
add_rect(sl, 0.3, 5.18, 12.73, 1.55, fill=NAVY)
tb(sl, "The flat WoLF-PHC line at ~9.1 bits/s/Hz across all jammer power levels is the paper's core claim.\n\n"
       "It occurs because η ≈ 0.13 does NOT change as jammer power increases — the jammer cannot force\n"
       "the agent to become more predictable. Unpredictability is an intrinsic algorithmic security property.",
   0.45, 5.25, 12.3, 1.4, fs=14, color=WHITE)

# ══════ SLIDE 15 – SINR CDF & Training Time ═══════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "Results: SINR CDF (Fig. 8) & Training Time (Fig. 9)",
    "QoS tail performance and computational cost analysis")
add_rect(sl, 0.3, 1.18, 6.2, 5.62, fill=LGRAY, line=BLUE, lw=Pt(1))
tb(sl, "Figure 8 — Empirical SINR CDF", 0.45, 1.22, 5.9, 0.38, fs=15, bold=True, color=NAVY)
tb(sl, "F(γ) = P(SINR ≤ γ)\n= fraction of (user, subcarrier, timestep)\n"
       "measurements with SINR ≤ γ\n\n"
       "Curve shifted RIGHT = higher SINR\n"
       "Steeper = more consistent\n"
       "10th pct (y=0.1) = worst-case tail\n\n"
       "Results:\n"
       "  WoLF-PHC:    10th pct ≈ +2 dB  ✓\n"
       "  AO Baseline: 10th pct ≈ −5 dB  △\n"
       "  Q-learning:  10th pct ≈ −25 dB ✗\n\n"
       "Q-learning has a FLAT PLATEAU near −25 to\n"
       "−15 dB: ~15% of timesteps the jammer fully\n"
       "exploited the agent → near-zero SINR.\n\n"
       "WoLF-PHC: NO plateau — CDF rises steeply.\n"
       "Steep shape confirms low SINR variance.",
   0.45, 1.65, 5.8, 5.1, fs=13, color=DGRAY)
add_rect(sl, 6.8, 1.18, 6.2, 5.62, fill=LGRAY, line=BLUE, lw=Pt(1))
tb(sl, "Figure 9 — Training Wall-Clock Time", 6.95, 1.22, 5.9, 0.38, fs=15, bold=True, color=NAVY)
tw = [2.3, 1.5, 2.0]
ch(sl,"Method",6.8,1.65,2.3,fs=12); ch(sl,"Time",9.1,1.65,1.5,fs=12); ch(sl,"Reason",10.6,1.65,2.0,fs=12)
for r, row in enumerate([("AO Baseline","2.5s","Matrix ops only"),
                           ("Q-learning","5.9s","Table + Bellman"),
                           ("Fast Q-learning","5.9s","Same as Q-learning"),
                           ("DQN","11.6s","NN fwd/bwd"),
                           ("WoLF-PHC","23.3s","3× eps + fuzzy + replay")]):
    tr(sl, row, 6.8, 2.03+r*0.38, tw, h=0.37,
       bg=RGBColor(0xD4,0xED,0xFF) if r==4 else (LGRAY if r%2==0 else WHITE), bold=(r==4))
tb(sl, "23.3s = one-time training cost.\n\n"
       "At inference: WoLF-PHC = fuzzy table\nlookup → same speed as Q-learning.\n\n"
       "The 4× training overhead justified by:\n"
       "  • 43% higher rate\n"
       "  • 75% higher SINR protection\n"
       "  • 14× lower deployment variance\n\n"
       "For a 6G base station, 23 seconds of\none-time training is negligible.",
   6.95, 3.88, 5.9, 2.85, fs=13, color=DGRAY)
tb(sl, "Training cost is paid once; deployment benefit is gained for every transmission.",
   0.35, 6.9, 12.6, 0.32, fs=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════ SLIDE 16 – Strong Points & Limitations ════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "Strong Points & Limitations",
    "Honest assessment of the proposed system")
tb(sl, "Strong Points:", 0.35, 1.18, 6.3, 0.38, fs=17, bold=True, color=GREEN)
for i, (title, body) in enumerate([
        ("Security-by-design",
         "WoLF-PHC's mixed strategy is INTRINSICALLY unpredictable. "
         "Jammer cannot elicit higher η regardless of power."),
        ("Adversarial robustness",
         "Flat ~9.1 bits/s/Hz as P_J goes 5→20 dBm while all baselines degrade significantly."),
        ("Reliable deployment",
         "14× lower variance across seeds (±0.18 vs ±2.52). Every seed, every channel."),
        ("Practical THz architecture",
         "SPDP-RIS + hybrid BF addresses real hardware constraints cost-effectively."),
        ("Tractable optimization",
         "RL+AO decomposes 256 continuous + discrete power without sacrificing quality.")]):
    y = 1.62 + i*1.02
    add_rect(sl, 0.3, y, 0.15, 0.75, fill=GREEN)
    tb(sl, title, 0.55, y+0.02, 3.0, 0.35, fs=13, bold=True, color=GREEN)
    tb(sl, body, 0.55, y+0.38, 5.8, 0.55, fs=12.5, color=DGRAY)
tb(sl, "Limitations & Future Work:", 7.0, 1.18, 6.0, 0.38, fs=17, bold=True, color=RED)
for i, (title, body) in enumerate([
        ("Single-cell only",
         "Multi-cell extension needs multi-agent RL and IRS coordination."),
        ("Perfect CSI assumed",
         "Imperfect/delayed CSI robustness needs investigation."),
        ("Near-field effects",
         "Large-aperture THz IRS in near-field violates plane-wave assumption."),
        ("Jammer channel partially known",
         "Full h_{J,k} knowledge may not be available in practice."),
        ("3× training budget",
         "Fair comparison would require equal compute budgets across methods.")]):
    y = 1.62 + i*1.02
    add_rect(sl, 7.0, y, 0.15, 0.75, fill=RED)
    tb(sl, title, 7.25, y+0.02, 3.0, 0.35, fs=13, bold=True, color=RED)
    tb(sl, body, 7.25, y+0.38, 5.7, 0.55, fs=12.5, color=DGRAY)

# ══════ SLIDE 17 – Conclusion ═════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
hdr(sl, "Conclusion", "Summary of contributions, results and outlook")
tb(sl, "What We Built:", 0.35, 1.18, 12.6, 0.38, fs=16, bold=True, color=NAVY)
for i, item in enumerate([
        "First complete THz IRS anti-jamming framework: 100 GHz, 10 GHz BW, M=256 IRS, N=64 BS, K=4 users, smart adaptive jammer",
        "SPDP-RIS with Q=64 TTD elements compensates beam squint across all M_sc=64 OFDM subcarriers",
        "Sub-connected hybrid BF with MVDR digital precoding — practical 8 RF-chain 100 GHz hardware solution",
        "Fuzzy WoLF-PHC: mixed strategy RL keeping η≈0.13, limits jammer boost to +2.3 dB vs +15–18 dB for deterministic methods",
        "Hybrid RL+AO: RL handles 42 discrete power actions, AO handles 256 continuous IRS phases (closed-form, 6 iterations)"]):
    tb(sl, f"▸  {item}", 0.35, 1.62+i*0.48, 12.6, 0.44, fs=13.5, color=DGRAY)
tb(sl, "Key Quantitative Results:", 0.35, 4.06, 12.6, 0.38, fs=16, bold=True, color=NAVY)
for i, (val, lbl, col) in enumerate([
        ("11.46\nbits/s/Hz", "System rate\n(vs 8.01 best baseline)", BLUE),
        ("78.6%", "SINR protection\n(vs 44.8%)", GREEN),
        ("±0.18 std", "Cross-seed variance\n(vs ±2.52 Q-learning)", NAVY),
        ("Flat ~9.1\nacross P_J", "Jammer immunity\n(all others degrade)", RED)]):
    x = 0.3 + i*3.18
    add_rect(sl, x, 4.5, 3.0, 1.3, fill=col)
    tb(sl, val, x+0.1, 4.55, 2.8, 0.58, fs=20, bold=True,
       color=WHITE, align=PP_ALIGN.CENTER)
    tb(sl, lbl, x+0.1, 5.12, 2.8, 0.62, fs=12,
       color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
tb(sl, "Future Directions:", 0.35, 5.95, 12.6, 0.38, fs=15, bold=True, color=NAVY)
tb(sl, "Multi-cell IRS coordination with multi-agent RL  •  Imperfect/delayed CSI robustness  •  "
       "Near-field THz IRS modelling  •  Hardware prototype at 100 GHz",
   0.35, 6.38, 12.6, 0.45, fs=13.5, color=DGRAY)
add_rect(sl, 0.3, 6.95, 12.73, 0.38, fill=NAVY)
tb(sl, "Thank you  |  Questions welcome  |  anuprabh@student.nitw.ac.in",
   0.4, 6.98, 12.4, 0.3, fs=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/anuprabh/Desktop/Secured Comm/IRS_AntiJamming_THz_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
